from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- 프레젠테이션 객체 생성 ---
prs = Presentation()

# --- 슬라이드 레이아웃 정의 ---
title_slide_layout = prs.slide_layouts[0] # 제목 슬라이드
title_and_content_layout = prs.slide_layouts[1] # 제목 및 내용 슬라이드
title_only_layout = prs.slide_layouts[5] # 제목만 있는 슬라이드

# --- 슬라이드 1: 제목 슬라이드 ---
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "YOLOv11n 모델 성능 개선"
subtitle.text = "내부 동작 원리 이해 및 실험 설계 보고\n\n발표자: 이진호\n2025. 09. 16."

# --- 슬라이드 2: 동작 원리 ---
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "후처리 하이퍼파라미터(Confidence & IoU) 동작 원리"

content_shape = slide.shapes.placeholders[1]
tf = content_shape.text_frame
tf.clear()

# Confidence Threshold 설명
p1 = tf.add_paragraph()
p1.text = "1. Confidence Threshold 란?"
p1.font.bold = True
p1.font.size = Pt(24)

p2 = tf.add_paragraph()
p2.text = "모델이 예측한 Bounding Box를 '객체가 맞다'고 최종 판단하는 최소한의 신뢰도(확률) 기준"
p2.level = 1
p2.font.size = Pt(20)

p3 = tf.add_paragraph()
p3.text = " (예시) Threshold = 0.5 → 신뢰도 50% 미만의 예측은 모두 무시"
p3.level = 2
p3.font.size = Pt(18)

tf.add_paragraph()

# NMS IoU Threshold 설명
p4 = tf.add_paragraph()
p4.text = "2. NMS IoU Threshold 란?"
p4.font.bold = True
p4.font.size = Pt(24)

p5 = tf.add_paragraph()
p5.text = "하나의 객체에 여러 Box가 중복으로 예측되었을 때, 이들을 하나의 Box로 합치는 기준 (NMS: Non-Maximum Suppression)"
p5.level = 1
p5.font.size = Pt(20)

p6 = tf.add_paragraph()
p6.text = "(예시) IoU > 0.7 → 70% 이상 겹치는 박스들은 같은 객체로 판단하여 하나로 통합"
p6.level = 2
p6.font.size = Pt(18)


# --- 슬라이드 3: 실험 계획 ---
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "성능 개선을 위한 실험 계획 (진행 중)"

# 테이블 추가
rows, cols = 3, 3
left = Inches(1.5)
top = Inches(2.0)
width = Inches(7.0)
height = Inches(2.5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 테이블 헤더
table.cell(0, 0).text = "실험 구분"
table.cell(0, 1).text = "개선 목표"
table.cell(0, 2).text = "파라미터 조건"

# 테이블 내용
table.cell(1, 0).text = "실험 1"
table.cell(1, 1).text = "중복 박스 문제 개선"
table.cell(1, 2).text = "IoU Threshold 조정 (0.7 → 0.5)"

table.cell(2, 0).text = "실험 2"
table.cell(2, 1).text = "재현율 개선"
table.cell(2, 2).text = "Confidence Threshold 조정 (0.25 → 0.1)"

# 테이블 서식 설정 (폰트, 정렬 등)
for row in table.rows:
    for cell in row.cells:
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # [최종 수정] import 대신 숫자 값을 직접 사용하여 세로 중앙 정렬
        cell.vertical_anchor = 3 # 3은 세로 중앙 정렬(Middle)을 의미하는 정수 값입니다.
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Malgun Gothic'
                run.font.size = Pt(18)

# 헤더 행 스타일 강조
for i in range(cols):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.text_frame.paragraphs[0].runs[0].font.bold = True

# --- 파일 저장 ---
file_name = "Week3_presentation.pptx"
prs.save(file_name)

print(f"'{file_name}' 파일이 성공적으로 생성되었습니다.")